"""
VERA Architecture Diagram — OpenVLA style  v6
==============================================
Left-to-right flow · white background · dashed VERA model border ·
individual token blocks · vertical merge-bar (no crossing arrows) ·
numbered steps ①②③④

Outputs:
  docs/VERA.png                   ← 200-dpi PNG for paper
  docs/VERA_architecture.pptx     ← editable 16×9 PowerPoint

Run:  python3 docs/make_vera_diagram_openvla.py
"""

import os
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.patches import FancyBboxPatch
import numpy as np

OUT_DIR = os.path.dirname(os.path.abspath(__file__))

# ── Canvas ────────────────────────────────────────────────────────────────────
FW, FH = 22.0, 10.5
DPI = 200
fig, ax = plt.subplots(figsize=(FW, FH))
ax.set_xlim(0, FW); ax.set_ylim(0, FH)
ax.axis("off"); ax.set_aspect("equal")
fig.patch.set_facecolor("white"); ax.set_facecolor("white")

# ── Colours ───────────────────────────────────────────────────────────────────
def hr(h):
    h = h.lstrip("#")
    return tuple(int(h[i:i+2], 16)/255 for i in (0, 2, 4))

C = dict(
    vis="#1565C0", ins="#2E7D32", act="#E65100",
    emb="#880E4F", his="#4527A0",
    clip="#37474F", sub="#4527A0",
    fuse="#0D47A1",
    dis="#BF360C",  reg="#1B5E20",
    new_="#C62828", gold="#F9A825",
    ice="#4DD0E1",  grn="#66BB6A",
    gray="#78909C", dark="#212121", mid="#546E7A",
    bord="#455A64", merge="#546E7A",
)

# ── Drawing helpers ───────────────────────────────────────────────────────────
def rbox(x, y, w, h, fc="none", ec="none", lw=1.5, dash=False,
         radius=0.08, alpha=1.0, z=3):
    ls = (0, (6, 3)) if dash else "solid"
    fc_ = hr(fc) if fc != "none" else (1, 1, 1, 0)
    ec_ = hr(ec) if ec != "none" else (0, 0, 0, 0)
    p = FancyBboxPatch((x, y), w, h,
                       boxstyle=f"round,pad={radius}",
                       facecolor=fc_, edgecolor=ec_,
                       linewidth=lw, linestyle=ls,
                       alpha=alpha, zorder=z, clip_on=False)
    ax.add_patch(p)

def txt(x, y, s, fs=9, c="#FFFFFF", ha="center", va="center",
        bold=False, italic=False, z=6):
    ax.text(x, y, s, ha=ha, va=va, fontsize=fs, color=c, zorder=z,
            fontweight="bold" if bold else "normal",
            fontstyle="italic" if italic else "normal", clip_on=False)

def arr(x1, y1, x2, y2, c="#78909C", lw=1.8, hw=0.16, style="arc3,rad=0"):
    ax.annotate("", xy=(x2, y2), xytext=(x1, y1),
        arrowprops=dict(
            arrowstyle=f"->,head_width={hw},head_length={hw*0.75}",
            color=c, lw=lw, connectionstyle=style),
        zorder=5, clip_on=False)

def hline(x1, x2, y, c="#78909C", lw=1.4):
    ax.plot([x1, x2], [y, y], color=c, lw=lw, zorder=4, clip_on=False)

def vline(x, y1, y2, c="#78909C", lw=1.4):
    ax.plot([x, x], [y1, y2], color=c, lw=lw, zorder=4, clip_on=False)

def badge(x, y, label, r=0.24):
    circ = plt.Circle((x, y), r, facecolor=hr("#1A237E"),
                      edgecolor=hr("#90CAF9"), lw=1.8, zorder=9, clip_on=False)
    ax.add_patch(circ)
    txt(x, y, label, fs=9.5, c="white", bold=True, z=10)

# ── Layout ────────────────────────────────────────────────────────────────────
STRH = 1.05          # input stream box height

# X zones (all in figure-inch coords)
INP_X,  INP_W  = 0.15, 2.72    # left input boxes

# Right output panel — set FIRST so model box fits inside figure
OUTW   = 2.05                          # output panel width
OUTX   = FW - OUTW - 0.18             # output panel left edge  (≈19.77)

# Model dashed box — right edge sits 0.22" before output panel
MBX    = INP_X + INP_W + 0.50         # ≈ 3.37
MBW    = OUTX - MBX - 0.22            # ≈ 16.18  (fits inside figure)
MBY,   MBH    = 0.55, FH - 1.05       # model box bottom, height
MBYT   = MBY + MBH                    # model box top
MB_R   = MBX + MBW                    # model right edge ≈ 19.55

# Inside model: columns  (verified: all widths + gaps = MBW - left/right pad)
ENC_X,  ENC_W  = MBX + 0.28, 1.88    # encoder column
MERGE_X = ENC_X + ENC_W + 0.12       # vertical merge bar left edge
MERGE_W = 0.16                        # merge bar width
MERGE_CX = MERGE_X + MERGE_W / 2
TOK_X   = MERGE_X + MERGE_W + 0.40   # token row left
TOK_W   = 4.40                        # token row width (tokens may overflow right)
# Fusion TF: from right of token area to HEAD_X
HEAD_W  = 2.00                        # policy head column width
HEAD_X  = MB_R - 0.18 - HEAD_W       # head column left edge
FTF_X   = TOK_X + TOK_W + 0.35       # fusion transformer left
FTF_W   = HEAD_X - FTF_X - 0.22      # fills the space between token area and heads

# Stream Y centres
YS = dict(vis=8.80, ins=7.25, act=5.70, emb=4.15, his=2.60)
STREAM_ORDER = ["vis", "ins", "act", "emb", "his"]
TOK_CY = (YS["vis"] + YS["his"]) / 2   # ≈ 5.70 = middle

# Fusion TF dimensions
FTF_Y0 = MBY + 0.22
FTF_H  = MBH - 0.44
FTF_CY = FTF_Y0 + FTF_H / 2

# Policy heads
DH_H = FTF_H * 0.38
RH_H = FTF_H * 0.38
DH_Y0 = FTF_Y0 + FTF_H - DH_H - 0.10
RH_Y0 = FTF_Y0 + 0.10
DH_MID = DH_Y0 + DH_H / 2
RH_MID = RH_Y0 + RH_H / 2

# ─────────────────────────────────────────────────────────────────────────────
#  TITLE + SUBTITLE
# ─────────────────────────────────────────────────────────────────────────────
txt(FW/2, FH - 0.26,
    "VERA: Vision · Embodied Knowledge · Reasoning · Action",
    fs=14, c=C["dark"], bold=True, z=10)
txt(FW/2, 0.22,
    "5-stream closed-loop robot policy  ·  6-layer Bidirectional LLaMA "
    "Fusion Transformer  ·  K=4 action chunking  ·  λ_align = 0.10",
    fs=8.5, c=C["mid"], z=10)

# ─────────────────────────────────────────────────────────────────────────────
#  DASHED VERA MODEL BORDER
# ─────────────────────────────────────────────────────────────────────────────
rbox(MBX, MBY, MBW, MBH, fc="none", ec=C["bord"],
     lw=2.0, dash=True, radius=0.05, z=1)
txt(MBX + 0.55, MBYT - 0.24, "VERA",
    fs=11, c=C["dark"], bold=True, ha="left", z=10)

# ─────────────────────────────────────────────────────────────────────────────
#  INPUT STREAM BOXES  (outside model)
# ─────────────────────────────────────────────────────────────────────────────
stream_info = [
    # (key, stream name, concise example)
    ("vis", "Vision",             "frames t, t−1, t−2"),
    ("ins", "Instruction",        '"push block to red region"'),
    ("act", "Action Narration",   '"I pushed the block left"'),
    ("emb", "Embodied Knowledge", '"block moved closer to goal"'),
    ("his", "History",            "(←,0.12)  (←,0.09)  (↑,0.04)  …"),
]
for key, title, sub in stream_info:
    cy = YS[key]; y0 = cy - STRH/2
    rbox(INP_X, y0, INP_W, STRH, fc=C[key], z=4)
    txt(INP_X + INP_W/2, cy + 0.24, title, fs=9.5, bold=True, z=6)
    txt(INP_X + INP_W/2, cy - 0.24, sub,   fs=7.5, c="#ECEFF1", italic=True, z=6)

# [NEW] badge on Embodied Knowledge
ey0 = YS["emb"] - STRH/2
rbox(INP_X + INP_W - 0.70, ey0 + STRH - 0.33, 0.60, 0.26,
     fc=C["new_"], radius=0.04, z=8)
txt(INP_X + INP_W - 0.40, ey0 + STRH - 0.205, "NEW", fs=7, bold=True, z=9)

# ─────────────────────────────────────────────────────────────────────────────
#  STEP ①  BADGE
# ─────────────────────────────────────────────────────────────────────────────
badge(ENC_X + ENC_W/2, MBYT - 0.38, "①")

# ─────────────────────────────────────────────────────────────────────────────
#  ENCODERS
# ─────────────────────────────────────────────────────────────────────────────
# CLIP ViT-B/32  (vision stream) — TRAINABLE, green border
VH = 1.32
rbox(ENC_X, YS["vis"] - VH/2, ENC_W, VH,
     fc=C["vis"], ec=C["grn"], lw=2.8, z=4)
txt(ENC_X + ENC_W/2, YS["vis"] + 0.36, "CLIP ViT-B/32",        fs=9, bold=True, z=6)
txt(ENC_X + ENC_W/2, YS["vis"] + 0.08, "✓ fine-tuned",         fs=8.5, c=C["grn"], bold=True, z=6)
txt(ENC_X + ENC_W/2, YS["vis"] - 0.26, "lr=3e-6 · 197 patches",fs=7.5, c="#C5CAE9", z=6)

# CLIP Text (frozen) — instruction, action narration, embodied knowledge
for key in ("ins", "act", "emb"):
    EH = 0.86
    cy = YS[key]
    rbox(ENC_X, cy - EH/2, ENC_W, EH,
         fc=C["clip"], ec=C["ice"], lw=2.0, dash=True, z=4)
    txt(ENC_X + ENC_W/2, cy + 0.19, "CLIP Text Encoder", fs=8.5, bold=True, z=6)
    txt(ENC_X + ENC_W/2, cy - 0.20, "❄ FROZEN  ·  1 token", fs=8, c=C["ice"], z=6)

# History encoder + causal sub-transformer
HEH = 1.12; cy = YS["his"]
HW2 = (ENC_W - 0.08) / 2
rbox(ENC_X, cy - HEH/2, HW2, HEH, fc=C["his"], z=4)
txt(ENC_X + HW2/2, cy + 0.24, "History",   fs=8.5, bold=True, z=6)
txt(ENC_X + HW2/2, cy - 0.20, "Encoder",   fs=8, z=6)
rbox(ENC_X + HW2 + 0.08, cy - HEH/2, HW2, HEH, fc=C["sub"], z=4)
txt(ENC_X + HW2 + 0.08 + HW2/2, cy + 0.24, "2L Causal",    fs=8.5, bold=True, z=6)
txt(ENC_X + HW2 + 0.08 + HW2/2, cy - 0.20, "LLaMA Sub-TF", fs=8, z=6)

# ─────────────────────────────────────────────────────────────────────────────
#  ARROWS: input boxes → encoder left edge  (colour-coded per stream)
# ─────────────────────────────────────────────────────────────────────────────
for key in STREAM_ORDER:
    cy = YS[key]
    arr(INP_X + INP_W, cy, ENC_X - 0.02, cy, c=C[key], lw=2.0, hw=0.15)

# ─────────────────────────────────────────────────────────────────────────────
#  VERTICAL MERGE BAR  (collects all encoder outputs cleanly)
# ─────────────────────────────────────────────────────────────────────────────
MERGE_Y0 = YS["his"] - 0.40    # bottom of merge bar
MERGE_Y1 = YS["vis"] + 0.40    # top of merge bar
rbox(MERGE_X, MERGE_Y0, MERGE_W, MERGE_Y1 - MERGE_Y0,
     fc=C["merge"], ec="none", lw=0, radius=0.04, alpha=0.85, z=5)

# Horizontal lines: encoder right edges → merge bar
for key in STREAM_ORDER:
    cy = YS[key]
    hline(ENC_X + ENC_W, MERGE_CX - 0.01, cy, c=C[key], lw=1.8)

# ─────────────────────────────────────────────────────────────────────────────
#  STEP ②  +  PROJECTION LABEL
# ─────────────────────────────────────────────────────────────────────────────
badge(MERGE_CX, MBYT - 0.38, "②")

# Projection label between merge bar and token row
PROJ_MID_X = (MERGE_X + MERGE_W + TOK_X) / 2
txt(PROJ_MID_X, MBYT - 0.64, "RMSNorm + Linear → d=256",
    fs=7.5, c=C["dark"], bold=True, z=6)
txt(PROJ_MID_X, MBYT - 0.90, "ViLT modality embeddings",
    fs=7, c=C["mid"], z=6)

# Arrow: merge bar → token row
arr(MERGE_X + MERGE_W, TOK_CY, TOK_X - 0.02, TOK_CY,
    c=C["gray"], lw=2.2, hw=0.20)

# ─────────────────────────────────────────────────────────────────────────────
#  TOKEN SEQUENCE ROW  (individual coloured blocks)
# ─────────────────────────────────────────────────────────────────────────────
badge(TOK_X + TOK_W/2, MBYT - 0.38, "③")

# Token definitions: (key, short_label, count, subgroup_label)
token_defs = [
    ("vis",  "V",      5, "patches"),
    ("ins",  "L_ins",  1, "instr."),
    ("act",  "L_act",  1, "action"),
    ("emb",  "L_emb",  1, "embodied"),
    ("his",  "h",      4, "history"),
    ("fuse", "CLS",    1, ""),
]

TBW = 0.44   # token block width
TBH = 0.56   # token block height
TG  = 0.05   # gap within group
TGG = 0.18   # extra gap between groups

# Compute total width and starting x
total_w = sum(n * TBW + max(n - 1, 0) * TG for _, _, n, _ in token_defs)
total_w += (len(token_defs) - 1) * TGG
start_x = TOK_X + (TOK_W - total_w) / 2
start_x = max(start_x, TOK_X + 0.10)

tok_cx = {}   # group centre-x for arrows
cur_x = start_x
for key, label, cnt, sublbl in token_defs:
    fc = C.get(key, "#1A237E")
    xs = []
    for i in range(cnt):
        rbox(cur_x, TOK_CY - TBH/2, TBW, TBH, fc=fc, z=7)
        if cnt == 1:
            lbl = label
        elif i == 0:
            lbl = f"{label}₁"
        elif i == cnt - 1:
            lbl = f"{label}_{cnt}"
        else:
            lbl = "…"
        txt(cur_x + TBW/2, TOK_CY, lbl, fs=7.5, bold=True, z=8)
        xs.append(cur_x + TBW/2)
        cur_x += TBW + TG
    tok_cx[key] = (min(xs) + max(xs)) / 2
    cur_x += TGG - TG   # replace last inner gap with larger inter-group gap

    # Sublabel below group
    if sublbl:
        g_mid = (min(xs) + max(xs)) / 2
        txt(g_mid, TOK_CY - TBH/2 - 0.20, sublbl, fs=7, c=C["dark"], z=6)

# [NEW] badge above L_emb token block
if "emb" in tok_cx:
    ex = tok_cx["emb"]
    rbox(ex - 0.22, TOK_CY + TBH/2 + 0.07, 0.44, 0.24,
         fc=C["new_"], radius=0.04, z=9)
    txt(ex, TOK_CY + TBH/2 + 0.19, "NEW", fs=6.5, bold=True, z=10)

# Token row label (top)
txt(TOK_X + TOK_W/2, TOK_CY + TBH/2 + 0.55,
    "token sequence (after projection)",
    fs=7.5, c=C["mid"], z=6)

# ─────────────────────────────────────────────────────────────────────────────
#  STEP ④  +  FUSION TRANSFORMER
# ─────────────────────────────────────────────────────────────────────────────
badge(FTF_X + FTF_W/2, MBYT - 0.38, "④")

# Arrow: token row → fusion TF
arr(TOK_X + TOK_W, TOK_CY, FTF_X - 0.02, TOK_CY,
    c=C["gold"], lw=2.8, hw=0.24)

# Main fusion TF body
rbox(FTF_X, FTF_Y0, FTF_W, FTF_H, fc=C["fuse"], ec="#7986CB", lw=2.0, z=3)

# Header bar inside TF
rbox(FTF_X + 0.10, FTF_Y0 + FTF_H - 0.92, FTF_W - 0.20, 0.82, fc="#1A237E", z=4)
txt(FTF_X + FTF_W/2, FTF_Y0 + FTF_H - 0.56,
    "LLaMA Fusion Transformer",
    fs=11.5, bold=True, z=6)
txt(FTF_X + FTF_W/2, FTF_Y0 + FTF_H - 0.80,
    "6 Layers  ·  8 Heads  ·  d_model = 256  ·  d_ff = 1024  ·  RMSNorm · RoPE · SwiGLU",
    fs=7.5, z=6)

# BIDIRECTIONAL banner
rbox(FTF_X + 0.10, FTF_Y0 + FTF_H - 1.50, FTF_W - 0.20, 0.46,
     fc=C["fuse"], ec=C["gold"], lw=2.0, z=4)
txt(FTF_X + FTF_W/2, FTF_Y0 + FTF_H - 1.27,
    "◈  BIDIRECTIONAL — Full Cross-Stream Attention — No Causal Mask  ◈",
    fs=9, c=C["gold"], bold=True, z=6)

# Token sequence label inside TF
txt(FTF_X + FTF_W/2, FTF_Y0 + FTF_H - 1.86,
    "[ L_ins | L_act | L_emb | V₁…V₁₉₇ | h₁…h₄ | CLS ]",
    fs=8.5, c="#C5CAE9", bold=True, z=6)

# 6 layer boxes (bottom of TF)
LBW = (FTF_W - 0.30) / 6
LBH = 0.50
LBY = FTF_Y0 + 0.18
for i in range(6):
    lx = FTF_X + 0.15 + i * (LBW + 0.02)
    fc2 = "#1565C0" if i % 2 == 0 else "#1976D2"
    rbox(lx, LBY, LBW - 0.02, LBH, fc=fc2, ec="#90CAF9", lw=0.8, z=5)
    txt(lx + (LBW - 0.02)/2, LBY + LBH/2, f"L{i+1}", fs=8, c="#E3F2FD", z=6)

# InfoNCE loss note (middle band inside TF)
INFONC_Y0 = LBY + LBH + 0.18
INFONC_H  = FTF_Y0 + FTF_H - 1.95 - INFONC_Y0
if INFONC_H > 0.8:
    rbox(FTF_X + 0.10, INFONC_Y0, FTF_W - 0.20, INFONC_H,
         fc="#0A1F6E", ec=C["gold"], lw=1.5, z=4)
    imid = INFONC_Y0 + INFONC_H / 2
    txt(FTF_X + FTF_W/2, imid + 0.40,
        "InfoNCE Alignment Loss  (training only)",
        fs=8.5, bold=True, c=C["gold"], z=6)
    txt(FTF_X + FTF_W/2, imid + 0.10,
        "InfoNCE( L_emb, L_act ‖ L_ins )  ·  reward-weighted  exp(5·r̃)",
        fs=8, c="#FFF8E1", z=6)
    txt(FTF_X + FTF_W/2, imid - 0.25,
        "λ_align = 0.10",
        fs=11, bold=True, c=C["gold"], z=6)

# ─────────────────────────────────────────────────────────────────────────────
#  POLICY HEADS  (inside model, right of fusion TF)
# ─────────────────────────────────────────────────────────────────────────────
# CLS arrow
arr(FTF_X + FTF_W, FTF_CY, HEAD_X - 0.02, FTF_CY,
    c=C["gold"], lw=2.5, hw=0.20)

# Discrete action head
rbox(HEAD_X, DH_Y0, HEAD_W, DH_H, fc=C["dis"], z=4)
txt(HEAD_X + HEAD_W/2, DH_Y0 + DH_H*0.74, "Discrete",      fs=9.5, bold=True, z=6)
txt(HEAD_X + HEAD_W/2, DH_Y0 + DH_H*0.54, "Action Head",   fs=9.5, bold=True, z=6)
txt(HEAD_X + HEAD_W/2, DH_Y0 + DH_H*0.34, "8-bin logits",  fs=8.5, z=6)
txt(HEAD_X + HEAD_W/2, DH_Y0 + DH_H*0.16, "× K=4 chunks",  fs=8,   z=6)

# Continuous regression head
rbox(HEAD_X, RH_Y0, HEAD_W, RH_H, fc=C["reg"], z=4)
txt(HEAD_X + HEAD_W/2, RH_Y0 + RH_H*0.74, "Regression",       fs=9.5, bold=True, z=6)
txt(HEAD_X + HEAD_W/2, RH_Y0 + RH_H*0.54, "Head",             fs=9.5, bold=True, z=6)
txt(HEAD_X + HEAD_W/2, RH_Y0 + RH_H*0.34, "RMSNorm + Tanh",   fs=8.5, z=6)
txt(HEAD_X + HEAD_W/2, RH_Y0 + RH_H*0.16, "[Δx, Δy]",         fs=9,   bold=True, z=6)

# ─────────────────────────────────────────────────────────────────────────────
#  RIGHT OUTPUT PANEL  (outside VERA model)
# ─────────────────────────────────────────────────────────────────────────────
arr(HEAD_X + HEAD_W, DH_MID, OUTX - 0.02, DH_MID,
    c="#FFCCBC", lw=2.0, hw=0.16)
arr(HEAD_X + HEAD_W, RH_MID, OUTX - 0.02, RH_MID,
    c="#C8E6C9", lw=2.0, hw=0.16)

OUTMID = OUTX + OUTW / 2

rbox(OUTX, DH_Y0, OUTW, DH_H, fc=C["dis"], z=4)
txt(OUTMID, DH_Y0 + DH_H*0.74, "Discrete Action",    fs=9, bold=True, z=6)
txt(OUTMID, DH_Y0 + DH_H*0.52, "8-bin × K=4",        fs=9, z=6)
txt(OUTMID, DH_Y0 + DH_H*0.30, "(π₀ / GR-1 style)",  fs=7.5, italic=True, c="#FFCCBC", z=6)

rbox(OUTX, RH_Y0, OUTW, RH_H, fc=C["reg"], z=4)
txt(OUTMID, RH_Y0 + RH_H*0.74, "Continuous Action",  fs=9, bold=True, z=6)
txt(OUTMID, RH_Y0 + RH_H*0.52, "[Δx, Δy]",           fs=9.5, bold=True, z=6)
txt(OUTMID, RH_Y0 + RH_H*0.30, "action_dim = 2",      fs=8, c="#C8E6C9", z=6)

# Closed-loop note (small text, no messy arc)
txt((FW/2), MBY - 0.22,
    "⟲  closed-loop: action narration (L_act) and embodied consequence (L_emb) "
    "re-injected as stream inputs at t+1",
    fs=7.5, c=C["mid"], z=6)

# ── Save PNG ──────────────────────────────────────────────────────────────────
png_out = os.path.join(OUT_DIR, "VERA.png")
plt.savefig(png_out, dpi=DPI, bbox_inches="tight", facecolor="white")
print(f"✓ Saved PNG → {png_out}")
plt.close()


# =============================================================================
#  POWERPOINT  (mirrors the matplotlib layout)
# =============================================================================
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from lxml import etree
from pptx.oxml.ns import qn

SW_P, SH_P = 16.0, 9.0   # slide size (inches)
prs = Presentation()
prs.slide_width  = Inches(SW_P)
prs.slide_height = Inches(SH_P)
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Scale: figure coords → slide inches
SXP = SW_P / FW
SYP = SH_P / FH

def prgb(h):
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

def _px(x): return x * SXP
def _py(y): return (FH - y) * SYP  # PPTX y=0 at top

def pbox(lft, top, w, h, fc, ec=None, lw=Pt(0), dash=False,
         lines=(), fsizes=(), bolds=(), tc="#FFFFFF"):
    shp = slide.shapes.add_shape(
        5, Inches(_px(lft)), Inches(_py(top + h)),
        Inches(w * SXP), Inches(h * SYP))
    try: shp.adjustments[0] = 0.04
    except Exception: pass
    shp.fill.solid()
    shp.fill.fore_color.rgb = prgb(fc)
    ln = shp.line
    if ec and ec != "none":
        ln.color.rgb = prgb(ec); ln.width = lw
        if dash:
            try: ln.dash_style = 4
            except Exception: pass
    else:
        ln.fill.background()
    tf = shp.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tc_rgb = prgb(tc)
    for i, (text, fs) in enumerate(zip(lines, fsizes or [9]*len(lines))):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = PP_ALIGN.CENTER
        run = para.add_run(); run.text = text
        run.font.size = Pt(fs)
        run.font.bold = (i in bolds)
        run.font.color.rgb = tc_rgb
    return shp

def ptxt(lft, top, w, h, text, fs=9, bold=False, color="#212121",
         align=PP_ALIGN.CENTER):
    txb = slide.shapes.add_textbox(
        Inches(_px(lft)), Inches(_py(top + h)),
        Inches(w * SXP), Inches(h * SYP))
    tf = txb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(fs); r.font.bold = bold
    r.font.color.rgb = prgb(color)

def parr(x1, y1, x2, y2, color="#78909C", lw=Pt(1.8)):
    conn = slide.shapes.add_connector(
        1, Inches(_px(x1)), Inches(_py(y1)),
           Inches(_px(x2)), Inches(_py(y2)))
    conn.line.color.rgb = prgb(color); conn.line.width = lw
    ln = conn.line._ln
    for tag in [qn("a:tailEnd"), qn("a:headEnd")]:
        for el in ln.findall(tag): ln.remove(el)
    tail = etree.SubElement(ln, qn("a:tailEnd")); tail.set("type", "none")
    head = etree.SubElement(ln, qn("a:headEnd"))
    head.set("type", "arrow"); head.set("w", "med"); head.set("len", "med")

# White background
bg = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(SW_P), Inches(SH_P))
bg.fill.solid(); bg.fill.fore_color.rgb = prgb("#FFFFFF")
bg.line.fill.background()

# Title
ptxt(0.15, FH - 0.55, FW - 0.30, 0.50,
     "VERA: Vision · Embodied Knowledge · Reasoning · Action",
     fs=16, bold=True, color="#212121")

# Model border
pbox(MBX, MBY, MBW, MBH, fc="#FAFAFA", ec=C["bord"], lw=Pt(2.0), dash=True)
ptxt(MBX + 0.40, MBY + MBH - 0.42, 1.80, 0.38,
     "VERA", fs=11, bold=True, color="#212121", align=PP_ALIGN.LEFT)

# Input stream boxes
for key, title, sub in stream_info:
    cy = YS[key]; y0 = cy - STRH/2
    sub_clean = sub.replace('“', '"').replace('”', '"')
    pbox(INP_X, y0, INP_W, STRH, fc=C[key],
         lines=[title, sub_clean], fsizes=[9, 7.5], bolds={0}, tc="#FFFFFF")

# [NEW] badge
ey0 = YS["emb"] - STRH/2
pbox(INP_X + INP_W - 0.70, ey0 + STRH - 0.33, 0.60, 0.26,
     fc=C["new_"], lines=["NEW"], fsizes=[7], bolds={0})

# Step badges (text labels)
for bx, by, bn in [
    (ENC_X + ENC_W/2,     MBYT - 0.38, "1"),
    (MERGE_CX,             MBYT - 0.38, "2"),
    (TOK_X + TOK_W/2,     MBYT - 0.38, "3"),
    (FTF_X + FTF_W/2,     MBYT - 0.38, "4"),
]:
    pbox(bx - 0.24, by - 0.24, 0.48, 0.48, fc="#1A237E", ec="#90CAF9", lw=Pt(1.8),
         lines=[bn], fsizes=[9], bolds={0})

# CLIP ViT
VH = 1.32
pbox(ENC_X, YS["vis"] - VH/2, ENC_W, VH,
     fc=C["vis"], ec=C["grn"], lw=Pt(2.8),
     lines=["CLIP ViT-B/32", "✓ fine-tuned", "lr=3e-6 · 197 patches"],
     fsizes=[9, 8.5, 7.5], bolds={0, 1})

# CLIP Text (frozen)
for key in ("ins", "act", "emb"):
    EH = 0.86; cy = YS[key]
    pbox(ENC_X, cy - EH/2, ENC_W, EH,
         fc=C["clip"], ec=C["ice"], lw=Pt(2.0), dash=True,
         lines=["CLIP Text Encoder", "❄ FROZEN  ·  1 token"],
         fsizes=[8.5, 8], bolds={0, 1}, tc="#FFFFFF")

# History
HEH = 1.12; cy = YS["his"]; HW2p = (ENC_W - 0.08) / 2
pbox(ENC_X, cy - HEH/2, HW2p, HEH,
     fc=C["his"], lines=["History", "Encoder"], fsizes=[8.5, 8], bolds={0})
pbox(ENC_X + HW2p + 0.08, cy - HEH/2, HW2p, HEH,
     fc=C["sub"], lines=["2L Causal", "LLaMA Sub-TF"], fsizes=[8.5, 8], bolds={0})

# Merge bar
pbox(MERGE_X, MERGE_Y0, MERGE_W, MERGE_Y1 - MERGE_Y0,
     fc=C["merge"], ec="none")

# Projection zone label
ptxt(PROJ_MID_X - 0.80, MBYT - 1.00, 1.60, 0.40,
     "RMSNorm + Linear\nd=256", fs=8, color="#212121")

# Fusion TF
pbox(FTF_X, FTF_Y0, FTF_W, FTF_H, fc=C["fuse"], ec="#7986CB", lw=Pt(2.0))
pbox(FTF_X + 0.10, FTF_Y0 + FTF_H - 0.92, FTF_W - 0.20, 0.82, fc="#1A237E",
     lines=["LLaMA Fusion Transformer",
            "6 Layers · 8 Heads · d=256 · d_ff=1024 · RMSNorm·RoPE·SwiGLU"],
     fsizes=[11.5, 7.5], bolds={0})
pbox(FTF_X + 0.10, FTF_Y0 + FTF_H - 1.50, FTF_W - 0.20, 0.46,
     fc=C["fuse"], ec=C["gold"], lw=Pt(2.0),
     lines=["◈  BIDIRECTIONAL — Full Cross-Stream Attention — No Causal Mask  ◈"],
     fsizes=[9], bolds={0}, tc=C["gold"])

LBW_P = (FTF_W - 0.30) / 6
LBH_P = 0.50
for i in range(6):
    lx = FTF_X + 0.15 + i * (LBW_P + 0.02)
    fc2 = "#1565C0" if i % 2 == 0 else "#1976D2"
    pbox(lx, LBY, LBW_P - 0.02, LBH_P,
         fc=fc2, lines=[f"L{i+1}"], fsizes=[8], bolds={0}, tc="#E3F2FD")

# Policy heads
pbox(HEAD_X, DH_Y0, HEAD_W, DH_H, fc=C["dis"],
     lines=["Discrete", "Action Head", "8-bin logits", "×K=4 chunks"],
     fsizes=[9.5, 9.5, 8.5, 8], bolds={0, 1})
pbox(HEAD_X, RH_Y0, HEAD_W, RH_H, fc=C["reg"],
     lines=["Regression", "Head", "[Δx, Δy]"],
     fsizes=[9.5, 9.5, 9], bolds={0, 1})

# Arrows
for key in STREAM_ORDER:
    cy = YS[key]
    parr(INP_X + INP_W, cy, MBX, cy, color=C[key], lw=Pt(2.0))
    parr(ENC_X + ENC_W, cy, MERGE_CX, cy, color=C[key], lw=Pt(1.8))
parr(MERGE_X + MERGE_W, TOK_CY, TOK_X, TOK_CY, color=C["gray"], lw=Pt(2.2))
parr(TOK_X + TOK_W, TOK_CY, FTF_X, TOK_CY, color=C["gold"], lw=Pt(2.8))
parr(FTF_X + FTF_W, FTF_CY, HEAD_X, FTF_CY, color=C["gold"], lw=Pt(2.5))
parr(HEAD_X + HEAD_W, DH_MID, OUTX, DH_MID, color="#FFCCBC", lw=Pt(2.0))
parr(HEAD_X + HEAD_W, RH_MID, OUTX, RH_MID, color="#C8E6C9", lw=Pt(2.0))

# Output boxes
pbox(OUTX, DH_Y0, OUTW, DH_H, fc=C["dis"],
     lines=["Discrete Action", "8-bin × K=4", "(π₀ style)"],
     fsizes=[9, 9, 7.5], bolds={0})
pbox(OUTX, RH_Y0, OUTW, RH_H, fc=C["reg"],
     lines=["Continuous", "[Δx, Δy]", "action_dim=2"],
     fsizes=[9, 9.5, 8], bolds={0})

# Subtitle
ptxt(0.15, 0.05, FW - 0.30, 0.32,
     "5-stream closed-loop robot policy  ·  6-layer Bidirectional LLaMA "
     "Fusion Transformer  ·  K=4 action chunking  ·  λ_align=0.10",
     fs=8, color="#546E7A")

pptx_out = os.path.join(OUT_DIR, "VERA_architecture.pptx")
prs.save(pptx_out)
print(f"✓ Saved PPTX → {pptx_out}")
