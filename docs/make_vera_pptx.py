"""
VERA Architecture Diagram  —  v4
Bottom-to-top layout (Open-VLA style):
  inputs at BOTTOM  →  arrows pointing UP  →  action outputs at TOP

Outputs:
  docs/VERA_architecture.pptx   ← editable 16×9 PowerPoint
  docs/VERA.png                  ← 180-dpi PNG for the paper

Run:  python3 docs/make_vera_pptx.py
"""

# ═══════════════════════════════════════════════════════════════════════════════
#  SHARED GEOMETRY  (used by both PPTX and PNG sections)
# ═══════════════════════════════════════════════════════════════════════════════
# Slide / figure size  (inches)
SW, SH = 16.0, 9.0          # slide  (16:9 widescreen)
FW, FH = 20.0, 11.5         # PNG figure (taller for full resolution)

# Five stream columns + right panel
MARGIN   = 0.22             # side margin
N_COLS   = 5
RIGHT_W  = 4.20             # InfoNCE panel width
GAP_COL  = 0.16             # horizontal gap between stream columns
STREAM_W = SW - 2*MARGIN - RIGHT_W - 0.22   # total width for 5 streams

COL_W = (STREAM_W - (N_COLS-1)*GAP_COL) / N_COLS   # per-stream column width

def col_x(c):
    """Left edge of stream column c (0-indexed)."""
    return MARGIN + c*(COL_W + GAP_COL)

RIGHT_X = MARGIN + STREAM_W + 0.22    # left edge of right panel

# Row heights  (bottom-to-top naming: r1=inputs, r6=outputs)
RH = dict(r1=1.00, r2=1.08, r3=0.38, r4=2.55, r5=0.52, r6=1.12)
RGAP = 0.20    # vertical gap between rows

TITLE_H    = 0.44
SUBTITLE_H = 0.28
LEGEND_H   = 0.42

# ── Compute row y-positions for PPTX (y=0 at TOP, increases downward) ─────────
#    Bottom-to-top: r6 (outputs) near top (small y), r1 (inputs) near bottom (large y)
_pptx_row = {}
_y = TITLE_H + SUBTITLE_H + 0.22          # start just below title
for row in ["r6","r5","r4","r3","r2","r1"]:   # outputs first (top) → inputs last (bottom)
    _pptx_row[row] = _y
    _y += RH[row] + RGAP

PPTX_LEGEND_TOP = _y + 0.05

# ── Compute row y-positions for matplotlib (y=0 at BOTTOM, increases upward) ──
#    Bottom-to-top: r1 (inputs) near bottom (small y), r6 (outputs) near top (large y)
PNG_SCALE = FH / SH          # scale factor to map slide inches → figure inches
_mpy = {}
_y2 = LEGEND_H + 0.20        # start just above legend
for row in ["r1","r2","r3","r4","r5","r6"]:   # inputs first (bottom) → outputs last (top)
    _mpy[row] = _y2 * PNG_SCALE
    _y2 += RH[row] + RGAP

PNG_TITLE_Y    = (_y2 + 0.18) * PNG_SCALE   # title y (top of figure)
PNG_LEGEND_Y   = 0.06 * PNG_SCALE

# ── Colour palette ─────────────────────────────────────────────────────────────
PAL = dict(
    vis  = "#1565C0",   ins  = "#2E7D32",   act  = "#E65100",
    emb  = "#AD1457",   his  = "#6A1B9A",
    vit  = "#1565C0",   clip = "#37474F",   sub  = "#4527A0",
    fuse = "#1A237E",   cls_ = "#0D47A1",
    dis  = "#BF360C",   reg  = "#1B5E20",   loss = "#F57F17",
    new  = "#B71C1C",   ice  = "#80DEEA",   grn  = "#A5D6A7",
    gold = "#FFD54F",   dark = "#1C1C1C",   mid  = "#555555",
    bg   = "#F4F6FB",   panel= "#ECEFF1",
)

STREAM_COLS = [PAL["vis"], PAL["ins"], PAL["act"], PAL["emb"], PAL["his"]]

# ═══════════════════════════════════════════════════════════════════════════════
#  POWERPOINT  SECTION
# ═══════════════════════════════════════════════════════════════════════════════
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from lxml import etree
from pptx.oxml.ns import qn
import os

def rgb(h):
    h = h.lstrip("#")
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

prs = Presentation()
prs.slide_width  = Inches(SW)
prs.slide_height = Inches(SH)
slide = prs.slides.add_slide(prs.slide_layouts[6])   # blank

# ── Background ────────────────────────────────────────────────────────────────
def _bg_rect():
    s = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(SW), Inches(SH))
    s.fill.solid(); s.fill.fore_color.rgb = rgb(PAL["bg"])
    s.line.fill.background()
_bg_rect()

# ── PPTX helpers ──────────────────────────────────────────────────────────────
def pbox(left, top, w, h, fc, ec=None, lw=Pt(0), dash=False, lines=(),
         fsizes=(), bolds=(), tc=None, alpha_fill=False):
    """Add a rounded rectangle with centred multi-line text."""
    shape = slide.shapes.add_shape(5, Inches(left), Inches(top),
                                   Inches(w), Inches(h))   # 5 = ROUNDED_RECTANGLE
    try: shape.adjustments[0] = 0.04
    except Exception: pass
    f = shape.fill
    if alpha_fill:
        f.background()
    else:
        f.solid(); f.fore_color.rgb = rgb(fc)
    ln = shape.line
    if ec:
        ln.color.rgb = rgb(ec); ln.width = lw
        if dash:
            try: ln.dash_style = 4
            except Exception: pass
    else:
        ln.fill.background()

    tf = shape.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    if not lines: return shape
    tc_rgb = rgb(tc) if tc else rgb("#FFFFFF")
    for i,(txt,fs) in enumerate(zip(lines, fsizes or [10]*len(lines))):
        para = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        para.alignment = PP_ALIGN.CENTER
        run = para.add_run(); run.text = txt
        run.font.size = Pt(fs); run.font.bold = (i in bolds)
        run.font.color.rgb = tc_rgb
    return shape

def ptxt(left, top, w, h, text, fs=9, bold=False, color=None, align=PP_ALIGN.CENTER):
    txb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
    tf  = txb.text_frame; tf.word_wrap = True
    p   = tf.paragraphs[0]; p.alignment = align
    r   = p.add_run(); r.text = text
    r.font.size = Pt(fs); r.font.bold = bold
    r.font.color.rgb = rgb(color) if color else rgb(PAL["dark"])

def parr(x1, y1, x2, y2, color, lw=Pt(2.2)):
    """Arrow from (x1,y1) to (x2,y2); arrowhead at (x2,y2)."""
    conn = slide.shapes.add_connector(1, Inches(x1), Inches(y1),
                                         Inches(x2), Inches(y2))
    conn.line.color.rgb = rgb(color); conn.line.width = lw
    ln = conn.line._ln
    # Remove any existing head/tail, add new ones
    for tag in [qn("a:tailEnd"), qn("a:headEnd")]:
        for el in ln.findall(tag): ln.remove(el)
    tail = etree.SubElement(ln, qn("a:tailEnd")); tail.set("type","none")
    head = etree.SubElement(ln, qn("a:headEnd"))
    head.set("type","arrow"); head.set("w","med"); head.set("len","med")

def row_cx(c): return col_x(c) + COL_W/2   # centre-x of stream column c
def row_mid(row): return _pptx_row[row] + RH[row]/2  # centre-y of row

# ── Title ─────────────────────────────────────────────────────────────────────
ptxt(MARGIN, 0.04, SW-2*MARGIN, TITLE_H,
     "VERA  —  Vision · Experience · Reasoning · Action",
     fs=20, bold=True, color=PAL["dark"])
ptxt(MARGIN, TITLE_H, SW-2*MARGIN, SUBTITLE_H,
     "5-stream closed-loop robot policy   ·   "
     "6-layer Bidirectional LLaMA Fusion Transformer   ·   "
     "K = 4 action chunking   ·   λ_align = 0.10",
     fs=9.5, color=PAL["mid"])

# ── Helper: upward arrows between adjacent rows ───────────────────────────────
def up_arrows_col(from_row, to_row, color=PAL["mid"], lw=Pt(1.8)):
    """Draw upward (in display sense) arrows for each stream column."""
    y1 = _pptx_row[from_row]               # top edge of lower row  (= arrow tail in PPTX)
    y2 = _pptx_row[to_row] + RH[to_row]   # bottom edge of upper row (= arrowhead)
    for c in range(N_COLS):
        parr(row_cx(c), y1, row_cx(c), y2, color, lw)

# ════════════════════════════════════════════════
# ROW 1  —  RAW INPUTS  (bottom of slide)
# ════════════════════════════════════════════════
r = "r1"; rt = _pptx_row[r]; rh = RH[r]

input_data = [
    (0, PAL["vis"], ["3 RGB Frames","224 × 224 × 3"],         [11,9],{0}),
    (1, PAL["ins"], ["Task Instruction",'"push block left"'],  [11,8],{0}),
    (2, PAL["act"], ["Prior Action  a_{t−1}","discrete bin [0–7]"], [10,8],{0}),
    (3, PAL["emb"], ["r_{t−1}  ·  Δd_{t−1}","reward + state delta"], [10,8],{0}),
    (4, PAL["his"], ["History Window","H = 4  (action, reward) pairs"], [10,8],{0}),
]
for c, fc, lines, fsz, bolds in input_data:
    pbox(col_x(c), rt, COL_W, rh, fc, lines=lines, fsizes=fsz, bolds=bolds)

# [NEW] badge on consequence column
nx = col_x(3)+COL_W-0.60; ny = rt+0.06
pbox(nx, ny, 0.55, 0.26, PAL["new"],
     lines=["NEW"], fsizes=[8], bolds={0})

up_arrows_col("r1","r2")

# ════════════════════════════════════════════════
# ROW 2  —  ENCODERS
# ════════════════════════════════════════════════
r = "r2"; rt = _pptx_row[r]; rh = RH[r]

# Col 0: CLIP ViT-B/32 — TRAINABLE (green border)
pbox(col_x(0), rt, COL_W, rh, PAL["vit"],
     ec=PAL["grn"], lw=Pt(3.0),
     lines=["CLIP ViT-B/32","✓  fine-tuned","lr = 3e-6","→ 197 patch tokens"],
     fsizes=[11,9,8.5,8], bolds={0,1})

# Cols 1-3: CLIP Text Encoder — FROZEN (dashed ice border)
for c in [1,2,3]:
    pbox(col_x(c), rt, COL_W, rh, PAL["clip"],
         ec=PAL["ice"], lw=Pt(3.0), dash=True,
         lines=["CLIP Text Encoder","❄  FROZEN","→ 1 token  (d=512)"],
         fsizes=[11,10,8], bolds={0,1})

# Col 4: History — two sub-boxes
hw2 = (COL_W-0.10)/2
pbox(col_x(4),         rt, hw2, rh, PAL["his"],
     lines=["History","Encoder","gate(a,r)"], fsizes=[10,9,8], bolds={0})
pbox(col_x(4)+hw2+0.10, rt, hw2, rh, PAL["sub"],
     lines=["2L Causal","LLaMA","Sub-TF"],   fsizes=[10,9,8], bolds={0})

up_arrows_col("r2","r3")

# ════════════════════════════════════════════════
# ROW 3  —  PROJECTION LABEL
# ════════════════════════════════════════════════
r = "r3"; rt = _pptx_row[r]; rh = RH[r]
pbox(MARGIN, rt, STREAM_W, rh, "#E0F2F1",
     ec="#00695C", lw=Pt(1.8),
     lines=["Independent  RMSNorm + Linear  →  d = 256   +   ViLT modality-type embedding per stream"],
     fsizes=[10.5], bolds={0}, tc="#004D40")

up_arrows_col("r3","r4")

# ════════════════════════════════════════════════
# ROW 4  —  LLAMA FUSION TRANSFORMER
# ════════════════════════════════════════════════
r = "r4"; rt = _pptx_row[r]; rh = RH[r]

# Outer body
pbox(MARGIN, rt, STREAM_W, rh, PAL["fuse"],
     ec="#7986CB", lw=Pt(2.5))

# Header bar
pbox(MARGIN+0.12, rt+0.10, STREAM_W-0.24, 0.48, "#283593",
     lines=["LLaMA Fusion Transformer   ·   6 Layers   ·   8 Heads   ·   d_model = 256   ·   d_ff = 1024   ·   RMSNorm · RoPE · SwiGLU"],
     fsizes=[10.5], bolds={0})

# BIDIRECTIONAL badge
pbox(MARGIN+0.12, rt+0.64, STREAM_W-0.24, 0.36, PAL["fuse"],
     ec=PAL["gold"], lw=Pt(2.5),
     lines=["◈   BIDIRECTIONAL  —  Full Attention  —  No Causal Mask   ◈"],
     fsizes=[10.5], bolds={0}, tc=PAL["gold"])

# Token sequence
pbox(MARGIN+0.12, rt+1.08, STREAM_W-0.24, 0.38, "#1A237E",
     lines=["[ t_ins  |  t_act  |  t_emb  |  V_patches  |  h₁ … h₄  |  CLS ]"],
     fsizes=[10], bolds={0}, tc="#C5CAE9")

# 6 layer boxes (stacked within fusion)
lw_each  = (STREAM_W-0.36)/6
lh_each  = 0.33
layer_top = rt + rh - 0.10 - lh_each   # bottom-most layer at bottom of fusion
for i in range(6):
    lx = MARGIN + 0.18 + i*(lw_each+0.02)
    # In PPTX (top-to-bottom), layer 1 is at the BOTTOM of the fusion box (largest y)
    pbox(lx, layer_top, lw_each, lh_each,
         "#1565C0" if i%2==0 else "#1976D2",
         ec="#90CAF9", lw=Pt(0.8),
         lines=[f"Layer {i+1}"], fsizes=[8], bolds={0})
    layer_top -= lh_each + 0.04   # stack upward inside the box

# ── Arrow fusion → CLS (single central arrow) ─────────────────────────────────
cx_all = MARGIN + STREAM_W/2
parr(cx_all, _pptx_row["r4"],             # top of r4 (= arrow tail)
     cx_all, _pptx_row["r5"]+RH["r5"],   # bottom of r5 (= arrowhead)
     PAL["gold"], lw=Pt(3.2))

# ════════════════════════════════════════════════
# ROW 5  —  CLS BAR
# ════════════════════════════════════════════════
r = "r5"; rt = _pptx_row[r]; rh = RH[r]
pbox(MARGIN, rt, STREAM_W, rh, PAL["cls_"],
     ec=PAL["gold"], lw=Pt(2.8),
     lines=["CLS  →  Policy Head Input"], fsizes=[13], bolds={0})

# ── Arrows CLS → two output heads ────────────────────────────────────────────
hw6 = (STREAM_W-0.20)/2
for offset in [hw6/2, hw6+0.20+hw6/2]:
    parr(cx_all, _pptx_row["r5"],
         MARGIN+offset, _pptx_row["r6"]+RH["r6"],
         "#BBDEFB", lw=Pt(2.5))

# ════════════════════════════════════════════════
# ROW 6  —  OUTPUT HEADS  (top of slide)
# ════════════════════════════════════════════════
r = "r6"; rt = _pptx_row[r]; rh = RH[r]
pbox(MARGIN, rt, hw6, rh, PAL["dis"],
     lines=["Discrete Action Head",
            "Expand-Compress → FC",
            "8-bin logits  ×  K = 4 chunks  (π₀ style)"],
     fsizes=[12,9.5,9.5], bolds={0})

pbox(MARGIN+hw6+0.20, rt, hw6, rh, PAL["reg"],
     lines=["Continuous Regression Head",
            "RMSNorm  +  Tanh",
            "Continuous action  [Δx,  Δy]   (action_dim = 2)"],
     fsizes=[12,9.5,9.5], bolds={0})

# ════════════════════════════════════════════════
# RIGHT PANEL  —  InfoNCE Loss
# ════════════════════════════════════════════════
# Spans from r6-top to r1-bottom (full content area)
loss_top = _pptx_row["r6"]
loss_bot = _pptx_row["r1"] + RH["r1"]
loss_h   = loss_bot - loss_top

pbox(RIGHT_X, loss_top, RIGHT_W, loss_h, PAL["loss"],
     ec="#FFFFFF", lw=Pt(2.0),
     lines=[
         "InfoNCE Alignment Loss",
         "(training only)",
         "",
         "InfoNCE( t_emb, t_act ‖ t_ins )",
         "reward-weighted · exp(5 · r̃)",
         "",
         "λ_align = 0.10",
         "",
         "Aligns consequence +",
         "action tokens toward",
         "the task instruction",
         "embedding",
     ],
     fsizes=[13,9,5,10.5,9.5,5,15,5,9.5,9.5,9.5,9.5],
     bolds={0,6})

# Closed-loop note below InfoNCE panel
cb_top = loss_bot + 0.10
pbox(RIGHT_X, cb_top, RIGHT_W, SH - cb_top - LEGEND_H - 0.15,
     "#ECEFF1", ec="#90A4AE", lw=Pt(1.2),
     lines=["⟲  Closed-loop feedback:",
            "action + consequence re-injected",
            "as stream inputs at  t+1"],
     fsizes=[10,9,9], bolds={0}, tc=PAL["dark"])

# ════════════════════════════════════════════════
# LEGEND  (bottom strip)
# ════════════════════════════════════════════════
lg_top  = PPTX_LEGEND_TOP
lg_h    = LEGEND_H
pbox(MARGIN, lg_top, SW-2*MARGIN, lg_h, "#ECEFF1",
     ec="#B0BEC5", lw=Pt(1.2), lines=[], fsizes=[])

leg_items = [
    (PAL["vis"],  None,        False, "Vision"),
    (PAL["ins"],  None,        False, "Instruction"),
    (PAL["act"],  None,        False, "Action Narration"),
    (PAL["emb"],  None,        False, "Embodied Consequence [NEW]"),
    (PAL["his"],  None,        False, "History"),
    (PAL["vit"],  PAL["grn"],  False, "CLIP ViT ✓ trainable (lr=3e-6)"),
    (PAL["clip"], PAL["ice"],  True,  "CLIP Text ❄ FROZEN"),
    (PAL["sub"],  None,        False, "2L Causal Sub-TF"),
    (PAL["fuse"], None,        False, "6L Bidir Fusion TF"),
    (PAL["dis"],  None,        False, "Discrete Head (8-bin×K=4)"),
    (PAL["reg"],  None,        False, "Regression Head [Δx,Δy]"),
    (PAL["loss"], None,        False, "InfoNCE Loss (λ=0.10)"),
]
item_w = (SW-2*MARGIN)/len(leg_items)
for j,(fc,ec_,dash,lbl) in enumerate(leg_items):
    lx = MARGIN + j*item_w
    pbox(lx+0.05, lg_top+0.06, 0.32, 0.26, fc,
         ec=ec_, lw=Pt(2.0) if ec_ else Pt(0), dash=dash)
    ptxt(lx+0.42, lg_top+0.06, item_w-0.44, 0.30, lbl, fs=7.5)

# ── Save PPTX ─────────────────────────────────────────────────────────────────
out_dir  = os.path.dirname(os.path.abspath(__file__))
pptx_out = os.path.join(out_dir, "VERA_architecture.pptx")
prs.save(pptx_out)
print(f"Saved PowerPoint → {pptx_out}")


# ═══════════════════════════════════════════════════════════════════════════════
#  MATPLOTLIB  PNG  SECTION   (bottom-to-top, y=0 at bottom)
# ═══════════════════════════════════════════════════════════════════════════════
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.patches import FancyBboxPatch
import numpy as np

fig, ax = plt.subplots(figsize=(FW, FH))
ax.set_xlim(0, FW);  ax.set_ylim(0, FH)
ax.set_aspect("equal"); ax.axis("off")
fig.patch.set_facecolor(PAL["bg"]); ax.set_facecolor(PAL["bg"])

# PNG scale: map slide coordinates → figure coordinates
SC = FW / SW   # horizontal scale
# For y: y=0 at BOTTOM in matplotlib; rows start at small y (bottom)

def h2c(h): return h.lstrip("#")
def c2r(h):
    h = h2c(h)
    return tuple(int(h[i:i+2],16)/255 for i in (0,2,4))

def mbox(x, y, w, h, fc, ec=None, lw=1.8, dash=False, alpha=1.0, z=3):
    """Draw a rounded box.  (x,y) = bottom-left corner in figure coords."""
    ls = (0,(5,3)) if dash else "solid"
    ec_ = c2r(ec) if ec else "none"
    p = FancyBboxPatch(
        (x*SC + 0.04, y*SC + 0.04), w*SC - 0.08, h*SC - 0.08,
        boxstyle="round,pad=0.04",
        facecolor=c2r(fc), edgecolor=ec_,
        linewidth=lw, linestyle=ls, alpha=alpha, zorder=z)
    ax.add_patch(p)

def mt(x, y, s, fs=9, c="#FFFFFF", ha="center", va="center",
       bold=False, italic=False, z=5):
    ax.text(x*SC, y*SC, s, ha=ha, va=va, fontsize=fs, color=c, zorder=z,
            fontweight="bold" if bold else "normal",
            fontstyle="italic" if italic else "normal")

def marr(x1,y1,x2,y2, c=PAL["mid"], lw=2.2):
    """Arrow from (x1,y1) to (x2,y2).  Arrowhead at (x2,y2)."""
    ax.annotate("", xy=(x2*SC, y2*SC), xytext=(x1*SC, y1*SC),
                arrowprops=dict(arrowstyle="->", color=c, lw=lw,
                                connectionstyle="arc3,rad=0"), zorder=4)

# Convenience: centre-x of stream column c
def mcx(c): return col_x(c) + COL_W/2

# Row bottom-left y in PNG (from _mpy which stores bottom edge of each row)
def rby(row): return _mpy[row] / SC    # undo SC to get slide-coord y
# Row top y
def rty(row): return rby(row) + RH[row]

# ── Legend strip (bottom) ─────────────────────────────────────────────────────
mbox(MARGIN, PNG_LEGEND_Y/SC, SW-2*MARGIN, LEGEND_H, "#ECEFF1",
     ec="#B0BEC5", lw=1.2, z=2)
leg_items_m = [
    (PAL["vis"],  None,        False, "Vision"),
    (PAL["ins"],  None,        False, "Instruction"),
    (PAL["act"],  None,        False, "Action Narration"),
    (PAL["emb"],  None,        False, "Embodied Consequence [NEW]"),
    (PAL["his"],  None,        False, "History"),
    (PAL["vit"],  PAL["grn"],  False, "CLIP ViT ✓ trainable (lr=3e-6)"),
    (PAL["clip"], PAL["ice"],  True,  "CLIP Text ❄ FROZEN"),
    (PAL["sub"],  None,        False, "2L Causal Sub-TF"),
    (PAL["fuse"], None,        False, "6L Bidir Fusion TF"),
    (PAL["dis"],  None,        False, "Discrete Head"),
    (PAL["reg"],  None,        False, "Regression Head"),
    (PAL["loss"], None,        False, "InfoNCE Loss (λ=0.10)"),
]
iw = (SW-2*MARGIN)/len(leg_items_m)
ly0 = PNG_LEGEND_Y/SC + 0.08
for j,(fc,ec,dash,lbl) in enumerate(leg_items_m):
    lx = MARGIN + j*iw
    mbox(lx+0.05, ly0, 0.28, 0.22, fc, ec=ec, lw=2.0, dash=dash, z=6)
    mt(lx+0.40, ly0+0.11, lbl, fs=7.5, c=PAL["dark"], ha="left", z=6)

# ── Title (top) ───────────────────────────────────────────────────────────────
title_y = PNG_TITLE_Y/SC
mt(SW/2, title_y + 0.28,
   "VERA  —  Vision · Experience · Reasoning · Action",
   fs=17, c=PAL["dark"], bold=True)
mt(SW/2, title_y,
   "5-stream closed-loop robot policy   ·   6-layer Bidirectional LLaMA Fusion Transformer   ·   K = 4   ·   λ_align = 0.10",
   fs=9, c=PAL["mid"])

# ════════════════════════════════════════════════
# ROW 1  —  INPUTS  (bottom)
# ════════════════════════════════════════════════
row = "r1"
input_info = [
    (0, PAL["vis"], "3 RGB Frames", "224 × 224 × 3"),
    (1, PAL["ins"], "Task Instruction", '"push block left"'),
    (2, PAL["act"], "Prior Action  a_{t−1}", "discrete bin [0–7]"),
    (3, PAL["emb"], "r_{t−1}  ·  Δd_{t−1}", "reward + state delta"),
    (4, PAL["his"], "History Window", "H = 4  (action, reward) pairs"),
]
for c, fc, l1, l2 in input_info:
    mbox(col_x(c), rby(row), COL_W, RH[row], fc)
    mt(mcx(c), rby(row)+RH[row]*0.64, l1, fs=9.5, bold=True)
    mt(mcx(c), rby(row)+RH[row]*0.28, l2, fs=8,  c="#E3F2FD")

# [NEW] badge
bx = col_x(3)+COL_W-0.60; by = rby(row)+RH[row]-0.30
mbox(bx, by, 0.55, 0.24, PAL["new"], z=7)
mt(bx+0.275, by+0.12, "NEW", fs=7.5, bold=True, z=8)

# Arrows r1 → r2
for c in range(N_COLS):
    marr(mcx(c), rty("r1"), mcx(c), rby("r2"))

# ════════════════════════════════════════════════
# ROW 2  —  ENCODERS
# ════════════════════════════════════════════════
row = "r2"
# CLIP ViT (trainable, green border)
mbox(col_x(0), rby(row), COL_W, RH[row], PAL["vit"], ec=PAL["grn"], lw=3.2)
mt(mcx(0), rby(row)+RH[row]*0.76, "CLIP ViT-B/32",      fs=10,  bold=True)
mt(mcx(0), rby(row)+RH[row]*0.52, "✓  fine-tuned",       fs=9,   c=PAL["grn"], bold=True)
mt(mcx(0), rby(row)+RH[row]*0.32, "lr = 3e-6",           fs=8.5, c=PAL["grn"])
mt(mcx(0), rby(row)+RH[row]*0.12, "197 patch tokens",    fs=8,   c="#C5CAE9")

# CLIP Text (frozen, dashed ice border)
for c in [1,2,3]:
    mbox(col_x(c), rby(row), COL_W, RH[row], PAL["clip"], ec=PAL["ice"], lw=3.2, dash=True)
    mt(mcx(c), rby(row)+RH[row]*0.76, "CLIP Text Encoder", fs=10,  bold=True)
    mt(mcx(c), rby(row)+RH[row]*0.50, "❄   FROZEN",         fs=9.5, c=PAL["ice"], bold=True)
    mt(mcx(c), rby(row)+RH[row]*0.22, "1 token  (d=512)",   fs=8,   c="#90A4AE")

# History: two sub-boxes
hw2 = (COL_W-0.10)/2
mbox(col_x(4),          rby(row), hw2, RH[row], PAL["his"])
mt(col_x(4)+hw2/2,      rby(row)+RH[row]*0.68, "History", fs=9.5, bold=True)
mt(col_x(4)+hw2/2,      rby(row)+RH[row]*0.40, "Encoder", fs=9)
mt(col_x(4)+hw2/2,      rby(row)+RH[row]*0.14, "gate(a,r)", fs=8, c="#D1C4E9")
mbox(col_x(4)+hw2+0.10, rby(row), hw2, RH[row], PAL["sub"])
mt(col_x(4)+hw2+0.10+hw2/2, rby(row)+RH[row]*0.72, "2L Causal", fs=9.5, bold=True)
mt(col_x(4)+hw2+0.10+hw2/2, rby(row)+RH[row]*0.46, "LLaMA",     fs=9.5, bold=True)
mt(col_x(4)+hw2+0.10+hw2/2, rby(row)+RH[row]*0.18, "Sub-TF",    fs=9)

# Arrows r2 → r3
for c in range(N_COLS):
    marr(mcx(c), rty("r2"), mcx(c), rby("r3"))

# ════════════════════════════════════════════════
# ROW 3  —  PROJECTION
# ════════════════════════════════════════════════
row = "r3"
mbox(MARGIN, rby(row), STREAM_W, RH[row], "#E0F2F1", ec="#00695C", lw=2.0)
mt(MARGIN+STREAM_W/2, rby(row)+RH[row]/2,
   "Independent  RMSNorm + Linear  →  d = 256   +   ViLT modality-type embedding per stream",
   fs=9.5, c="#004D40", bold=True)

# Arrows r3 → r4
for c in range(N_COLS):
    marr(mcx(c), rty("r3"), mcx(c), rby("r4"))

# ════════════════════════════════════════════════
# ROW 4  —  LLAMA FUSION TRANSFORMER
# ════════════════════════════════════════════════
row = "r4"
r4b = rby(row); r4h = RH[row]
mbox(MARGIN, r4b, STREAM_W, r4h, PAL["fuse"], ec="#7986CB", lw=2.8)

# Header bar (at TOP of fusion box in display = highest y)
mbox(MARGIN+0.12, r4b+r4h-0.54, STREAM_W-0.24, 0.48, "#283593")
mt(MARGIN+STREAM_W/2, r4b+r4h-0.30,
   "LLaMA Fusion Transformer   ·   6 Layers   ·   8 Heads   ·   d_model = 256   ·   d_ff = 1024   ·   RMSNorm · RoPE · SwiGLU",
   fs=9.5, bold=True)

# BIDIRECTIONAL banner
mbox(MARGIN+0.12, r4b+r4h-0.98, STREAM_W-0.24, 0.38, PAL["fuse"], ec=PAL["gold"], lw=2.5)
mt(MARGIN+STREAM_W/2, r4b+r4h-0.79,
   "◈   BIDIRECTIONAL  —  Full Attention  —  No Causal Mask   ◈",
   fs=9.5, c=PAL["gold"], bold=True)

# Token sequence
mt(MARGIN+STREAM_W/2, r4b+r4h-1.38,
   "[ t_ins  |  t_act  |  t_emb  |  V_patches  |  h₁ … h₄  |  CLS ]",
   fs=9.5, c="#C5CAE9", bold=True)

# 6 layer boxes (stacked at BOTTOM of fusion box)
lw_e  = (STREAM_W-0.36)/6
lh_e  = 0.27
lgap  = 0.06
l_bot = r4b + 0.12   # bottom of layer stack
for i in range(6):
    lx_e = MARGIN + 0.18 + i*(lw_e+0.02)
    mbox(lx_e, l_bot, lw_e, lh_e,
         "#1565C0" if i%2==0 else "#1976D2",
         ec="#90CAF9", lw=0.9, z=5)
    mt(lx_e+lw_e/2, l_bot+lh_e/2, f"Layer {i+1}", fs=7.5, c="#E3F2FD", z=6)

# Arrow fusion → CLS
cx_all = MARGIN + STREAM_W/2
marr(cx_all, rty("r4"), cx_all, rby("r5"), c=PAL["gold"], lw=3.2)

# ════════════════════════════════════════════════
# ROW 5  —  CLS BAR
# ════════════════════════════════════════════════
row = "r5"
mbox(MARGIN, rby(row), STREAM_W, RH[row], PAL["cls_"], ec=PAL["gold"], lw=3.0)
mt(MARGIN+STREAM_W/2, rby(row)+RH[row]/2,
   "CLS  →  Policy Head Input", fs=13, bold=True)

# Arrows CLS → output heads
hw6 = (STREAM_W-0.20)/2
for offset in [hw6/2, hw6+0.20+hw6/2]:
    marr(cx_all, rty("r5"), MARGIN+offset, rby("r6"), c="#BBDEFB", lw=2.5)

# ════════════════════════════════════════════════
# ROW 6  —  OUTPUT HEADS  (top)
# ════════════════════════════════════════════════
row = "r6"
# Discrete
mbox(MARGIN, rby(row), hw6, RH[row], PAL["dis"])
mt(MARGIN+hw6/2, rby(row)+RH[row]*0.76, "Discrete Action Head",        fs=11, bold=True)
mt(MARGIN+hw6/2, rby(row)+RH[row]*0.50, "Expand-Compress  →  FC",      fs=9,  c="#FFCCBC")
mt(MARGIN+hw6/2, rby(row)+RH[row]*0.22, "8-bin logits  ×  K = 4 chunks", fs=9, bold=True)

# Regression
rx = MARGIN+hw6+0.20
mbox(rx, rby(row), hw6, RH[row], PAL["reg"])
mt(rx+hw6/2, rby(row)+RH[row]*0.76, "Continuous Regression Head",     fs=11, bold=True)
mt(rx+hw6/2, rby(row)+RH[row]*0.50, "RMSNorm  +  Tanh",               fs=9,  c="#C8E6C9")
mt(rx+hw6/2, rby(row)+RH[row]*0.22, "Continuous action  [Δx,  Δy]",   fs=9,  bold=True)

# ════════════════════════════════════════════════
# RIGHT PANEL  —  InfoNCE Loss
# ════════════════════════════════════════════════
r_top  = rby("r1")
r_high = rty("r6")
r_h    = r_high - r_top
mbox(RIGHT_X, r_top, RIGHT_W, r_h, PAL["loss"])
mid_loss = r_top + r_h/2
mt(RIGHT_X+RIGHT_W/2, r_high - 0.42, "InfoNCE Alignment Loss", fs=12, bold=True)
mt(RIGHT_X+RIGHT_W/2, r_high - 0.78, "(training only)", fs=9,  italic=True)
mt(RIGHT_X+RIGHT_W/2, mid_loss+0.28, "InfoNCE( t_emb, t_act ‖ t_ins )", fs=10, bold=True)
mt(RIGHT_X+RIGHT_W/2, mid_loss-0.06, "reward-weighted · exp(5 · r̃)", fs=9,  c="#FFF8E1")
mt(RIGHT_X+RIGHT_W/2, mid_loss-0.46, "λ_align  =  0.10", fs=14, bold=True)
mt(RIGHT_X+RIGHT_W/2, r_top+0.60,    "Aligns consequence + action", fs=9,  c="#FFF8E1")
mt(RIGHT_X+RIGHT_W/2, r_top+0.30,    "toward task instruction", fs=9,  c="#FFF8E1")

# ── Save PNG ──────────────────────────────────────────────────────────────────
png_out = os.path.join(out_dir, "VERA.png")
plt.savefig(png_out, dpi=180, bbox_inches="tight", facecolor=PAL["bg"])
print(f"Saved PNG       → {png_out}")
